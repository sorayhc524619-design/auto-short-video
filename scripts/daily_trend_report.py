"""
scripts/daily_trend_report.py
毎日のUS BGM YouTube トレンドレポートを PowerPoint + Excel で生成。

実行: python scripts/daily_trend_report.py
出力: reports/YYYYMMDD_bgm_trend.pptx
     reports/YYYYMMDD_bgm_trend.xlsx

ClaudeにWeb調査 + 共通点分析 + 自チャンネル向け推奨を依頼。
Windowsタスクスケジューラから毎日 00:00 起動を想定。
"""

import json
import logging
import os
import sys
from datetime import datetime, timedelta
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
import config

import anthropic
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
)
logger = logging.getLogger(__name__)


CLAUDE_RESEARCH_PROMPT = """You are a YouTube trend analyst specializing in the US Background Music (BGM) / Sleep Music / Lo-Fi / Ambient niche.

Today's date: {today}
Target date for trend analysis: {target_date} ({weekday})

TASK: Output a JSON report covering the 10 currently-trending or recently-viral BGM videos in the US market (sleep music, lo-fi, ambient, cozy ambience, rain sounds, fireplace, study music). Focus on what would have been driving views in the last 24-72 hours.

Important constraints:
- Real, verifiable channel names and video title patterns (not hallucinated specific URLs)
- For each video give the title pattern and clearly mark whether it is an actual recent video you know about, or a representative pattern from the channel
- Common patterns must be specific and actionable (not generic platitudes)
- Each insight should suggest a concrete tweak for a small new channel to apply today

Output ONLY valid JSON, no markdown fences, in this exact shape:

{{
  "report_date": "{target_date}",
  "summary": "2-3 sentence executive summary of what is driving the US BGM niche right now",
  "videos": [
    {{
      "rank": 1,
      "title_pattern": "Cozy Rainy Cabin - 8 Hours of Rain Sounds for Sleep",
      "channel": "Cozy Cabin Ambience",
      "channel_size_band": "100K-1M",
      "theme": "rainy cabin",
      "duration_hours": 8,
      "visual_style": "static cabin window with rain effect",
      "music_style": "rain sounds with soft piano",
      "thumbnail_keywords": "warm window, dark exterior, rain droplets",
      "why_working": "specific 1-sentence reason this format wins",
      "key_takeaway": "actionable lesson for a new BGM channel"
    }}
    // ... 10 entries total
  ],
  "common_patterns": {{
    "title_patterns": ["pattern 1", "pattern 2", "pattern 3", "pattern 4", "pattern 5"],
    "visual_patterns": ["pattern 1", "..."],
    "audio_patterns": ["pattern 1", "..."],
    "duration_distribution": "e.g. 80% are 8-10 hours, 20% are 1-3 hours",
    "thumbnail_color_palette": "describe the dominant colors",
    "emoji_in_titles": ["🌧️", "❄️", "🔥", "🌙"],
    "seo_keywords_top_10": ["sleep music", "rain sounds", "..."]
  }},
  "today_recommendation": {{
    "theme_name": "english_name_for_today_video",
    "title": "exact YouTube title to use today (60-90 chars, English, includes hours and target use)",
    "description_first_paragraph": "first 2-3 sentences of description",
    "music_prompts_for_suno": [
      "5 Suno prompts, each ~40 words",
      "...",
      "...",
      "...",
      "..."
    ],
    "pixabay_search_keywords": ["english query 1", "english query 2", "english query 3"],
    "ambient_sound": "rain or fireplace or wind or none",
    "rationale": "why this theme today specifically"
  }}
}}
"""


def call_claude_for_report(target_date: datetime) -> dict:
    if not config.CLAUDE_API_KEY:
        raise RuntimeError("CLAUDE_API_KEY が設定されていません")
    client = anthropic.Anthropic(api_key=config.CLAUDE_API_KEY)
    prompt = CLAUDE_RESEARCH_PROMPT.format(
        today=datetime.now().strftime("%Y-%m-%d"),
        target_date=target_date.strftime("%Y-%m-%d"),
        weekday=target_date.strftime("%A"),
    )
    logger.info("Claudeにトレンド分析を依頼中...")
    msg = client.messages.create(
        model=config.CLAUDE_MODEL,
        max_tokens=8000,
        messages=[{"role": "user", "content": prompt}],
    )
    text = msg.content[0].text.strip()
    if text.startswith("```"):
        text = "\n".join(text.split("\n")[1:-1])
    data = json.loads(text)
    logger.info(f"レポート取得: {len(data.get('videos', []))} 件")
    return data


# ===== Excel =====

HEADER_FILL = PatternFill("solid", fgColor="1F2A44")
HEADER_FONT = Font(name="Calibri", size=11, color="FFFFFF", bold=True)
CELL_FONT = Font(name="Calibri", size=10)
ALIGN_TOP = Alignment(wrap_text=True, vertical="top", horizontal="left")


def write_excel(report: dict, output: Path):
    wb = Workbook()

    # --- Sheet 1: Top 10 Videos ---
    ws = wb.active
    ws.title = "Top 10 Videos"
    headers = [
        "Rank", "Title Pattern", "Channel", "Channel Size", "Theme",
        "Duration (h)", "Visual Style", "Music Style",
        "Thumbnail Keywords", "Why Working", "Key Takeaway",
    ]
    ws.append(headers)
    for c in range(1, len(headers) + 1):
        cell = ws.cell(row=1, column=c)
        cell.fill = HEADER_FILL
        cell.font = HEADER_FONT
        cell.alignment = ALIGN_TOP

    for v in report.get("videos", []):
        ws.append([
            v.get("rank"),
            v.get("title_pattern"),
            v.get("channel"),
            v.get("channel_size_band"),
            v.get("theme"),
            v.get("duration_hours"),
            v.get("visual_style"),
            v.get("music_style"),
            v.get("thumbnail_keywords"),
            v.get("why_working"),
            v.get("key_takeaway"),
        ])
    for c in range(1, len(headers) + 1):
        for r in range(2, ws.max_row + 1):
            ws.cell(row=r, column=c).font = CELL_FONT
            ws.cell(row=r, column=c).alignment = ALIGN_TOP
    widths = [6, 42, 22, 14, 18, 12, 28, 28, 28, 36, 36]
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w
    ws.freeze_panes = "A2"

    # --- Sheet 2: Common Patterns ---
    ws2 = wb.create_sheet("Common Patterns")
    cp = report.get("common_patterns", {})
    rows = [
        ("Category", "Items"),
        ("Title Patterns", "\n".join(cp.get("title_patterns", []))),
        ("Visual Patterns", "\n".join(cp.get("visual_patterns", []))),
        ("Audio Patterns", "\n".join(cp.get("audio_patterns", []))),
        ("Duration Distribution", cp.get("duration_distribution", "")),
        ("Thumbnail Color Palette", cp.get("thumbnail_color_palette", "")),
        ("Emoji In Titles", " ".join(cp.get("emoji_in_titles", []))),
        ("Top SEO Keywords", ", ".join(cp.get("seo_keywords_top_10", []))),
    ]
    for r, (k, v) in enumerate(rows, 1):
        ws2.cell(row=r, column=1, value=k)
        ws2.cell(row=r, column=2, value=v)
        ws2.cell(row=r, column=1).alignment = ALIGN_TOP
        ws2.cell(row=r, column=2).alignment = ALIGN_TOP
    ws2.cell(row=1, column=1).fill = HEADER_FILL
    ws2.cell(row=1, column=1).font = HEADER_FONT
    ws2.cell(row=1, column=2).fill = HEADER_FILL
    ws2.cell(row=1, column=2).font = HEADER_FONT
    ws2.column_dimensions["A"].width = 26
    ws2.column_dimensions["B"].width = 85

    # --- Sheet 3: Today's Recommendation ---
    ws3 = wb.create_sheet("Today's Action Plan")
    rec = report.get("today_recommendation", {})
    rows3 = [
        ("Field", "Value"),
        ("Theme Name", rec.get("theme_name", "")),
        ("Title", rec.get("title", "")),
        ("Description (intro)", rec.get("description_first_paragraph", "")),
        ("Ambient Sound", rec.get("ambient_sound", "")),
        ("Pixabay Keywords", ", ".join(rec.get("pixabay_search_keywords", []))),
        ("Rationale", rec.get("rationale", "")),
    ]
    for r, (k, v) in enumerate(rows3, 1):
        ws3.cell(row=r, column=1, value=k)
        ws3.cell(row=r, column=2, value=v)
        ws3.cell(row=r, column=1).alignment = ALIGN_TOP
        ws3.cell(row=r, column=2).alignment = ALIGN_TOP

    # Suno prompts: separate block
    ws3.append([])
    ws3.append(["Suno Prompts", ""])
    ws3.cell(row=ws3.max_row, column=1).fill = HEADER_FILL
    ws3.cell(row=ws3.max_row, column=1).font = HEADER_FONT
    ws3.cell(row=ws3.max_row, column=2).fill = HEADER_FILL
    ws3.cell(row=ws3.max_row, column=2).font = HEADER_FONT
    for i, p in enumerate(rec.get("music_prompts_for_suno", []), 1):
        ws3.append([f"Track {i}", p])
    ws3.cell(row=1, column=1).fill = HEADER_FILL
    ws3.cell(row=1, column=1).font = HEADER_FONT
    ws3.cell(row=1, column=2).fill = HEADER_FILL
    ws3.cell(row=1, column=2).font = HEADER_FONT
    ws3.column_dimensions["A"].width = 22
    ws3.column_dimensions["B"].width = 95
    for r in range(1, ws3.max_row + 1):
        ws3.row_dimensions[r].height = max(20, ws3.row_dimensions[r].height or 0)

    wb.save(output)
    logger.info(f"Excel保存: {output}")


# ===== PowerPoint =====

NAVY = RGBColor(0x1F, 0x2A, 0x44)
AMBER = RGBColor(0xFF, 0xB3, 0x47)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x55, 0x55)


def _add_title_slide(prs, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sb = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.3), Inches(1.5))
    sp = sb.text_frame.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(20)
    sp.font.color.rgb = AMBER


def _add_section_header(prs, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12.3), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = AMBER


def _add_content_slide(prs, title: str, body_lines: list):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # body
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.0))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, tuple):
            label, value = line
            para.text = f"{label}: "
            para.font.bold = True
            para.font.size = Pt(14)
            para.font.color.rgb = NAVY
            run = para.add_run()
            run.text = str(value or "")
            run.font.bold = False
            run.font.size = Pt(14)
            run.font.color.rgb = GRAY
        else:
            para.text = f"• {line}"
            para.font.size = Pt(14)
            para.font.color.rgb = GRAY


def write_pptx(report: dict, output: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    target_date = report.get("report_date", "")

    _add_title_slide(
        prs,
        "US BGM Trend Report",
        f"Daily Snapshot for {target_date}",
    )

    _add_content_slide(
        prs,
        "Executive Summary",
        [report.get("summary", "")],
    )

    # 10 videos
    _add_section_header(prs, "Top 10 Trending Videos")
    for v in report.get("videos", []):
        _add_content_slide(
            prs,
            f"#{v.get('rank')}  {v.get('title_pattern', '')[:60]}",
            [
                ("Channel", v.get("channel")),
                ("Channel Size", v.get("channel_size_band")),
                ("Theme", v.get("theme")),
                ("Duration", f"{v.get('duration_hours')} hours"),
                ("Visual", v.get("visual_style")),
                ("Music", v.get("music_style")),
                ("Thumbnail", v.get("thumbnail_keywords")),
                ("Why it works", v.get("why_working")),
                ("Takeaway", v.get("key_takeaway")),
            ],
        )

    # patterns
    _add_section_header(prs, "Common Patterns")
    cp = report.get("common_patterns", {})
    _add_content_slide(
        prs,
        "Title Patterns",
        cp.get("title_patterns", []),
    )
    _add_content_slide(
        prs,
        "Visual Patterns",
        cp.get("visual_patterns", []),
    )
    _add_content_slide(
        prs,
        "Audio Patterns",
        cp.get("audio_patterns", []),
    )
    _add_content_slide(
        prs,
        "Meta Patterns",
        [
            ("Duration Distribution", cp.get("duration_distribution")),
            ("Thumbnail Color Palette", cp.get("thumbnail_color_palette")),
            ("Emoji in Titles", " ".join(cp.get("emoji_in_titles", []))),
            ("Top SEO Keywords", ", ".join(cp.get("seo_keywords_top_10", []))),
        ],
    )

    # today's recommendation
    _add_section_header(prs, "Today's Action Plan")
    rec = report.get("today_recommendation", {})
    _add_content_slide(
        prs,
        "Recommended Theme for Today",
        [
            ("Theme", rec.get("theme_name")),
            ("Title", rec.get("title")),
            ("Description Opener", rec.get("description_first_paragraph")),
            ("Ambient Sound", rec.get("ambient_sound")),
            ("Pixabay Keywords", ", ".join(rec.get("pixabay_search_keywords", []))),
            ("Why this today", rec.get("rationale")),
        ],
    )

    _add_content_slide(
        prs,
        "Suno Prompts (Copy/Paste)",
        [
            f"Track {i+1}: {p}"
            for i, p in enumerate(rec.get("music_prompts_for_suno", [])[:5])
        ],
    )

    prs.save(output)
    logger.info(f"PowerPoint保存: {output}")


# ===== Main =====


def main():
    target = datetime.now() - timedelta(days=1)  # 「昨日」
    label = target.strftime("%Y%m%d")

    reports_dir = config.BASE_DIR / "reports"
    reports_dir.mkdir(exist_ok=True)

    pptx_path = reports_dir / f"{label}_bgm_trend.pptx"
    xlsx_path = reports_dir / f"{label}_bgm_trend.xlsx"
    json_path = reports_dir / f"{label}_bgm_trend.json"

    report = call_claude_for_report(target)
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(report, f, ensure_ascii=False, indent=2)

    write_excel(report, xlsx_path)
    write_pptx(report, pptx_path)

    logger.info("=" * 60)
    logger.info(f"完了:")
    logger.info(f"  {pptx_path}")
    logger.info(f"  {xlsx_path}")
    logger.info(f"  {json_path}")
    logger.info("=" * 60)


if __name__ == "__main__":
    main()
